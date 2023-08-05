#!/usr/bin/env python
# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0

import argparse
import boto3
import os
import traceback

from os import listdir
from os.path import isfile, join
from pathlib import Path
from botocore.exceptions import ClientError
from pptx import Presentation
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.text import MSO_AUTO_SIZE
from datetime import date
from pptx.util import Pt

LANGUAGE_CODE_TO_LANGUAGE_ID = {
    """
    Dict that maps Amazon Translate language code to MSO_LANGUAGE_ID enum value.
    - Amazon Translate language codes: https://docs.aws.amazon.com/translate/latest/dg/what-is.html#what-is-languages
    - python-pptx MSO_LANGUAGE_ID enum: https://python-pptx.readthedocs.io/en/latest/api/enum/MsoLanguageId.html
    python-pptx doesn't support:
        - Azerbaijani (az)
        - Persian (fa)
        - Dari (fa-AF)
        - Tagalog (tl)
    """
    'af': MSO_LANGUAGE_ID.AFRIKAANS,
    'am': MSO_LANGUAGE_ID.AMHARIC,
    'ar': MSO_LANGUAGE_ID.ARABIC,
    'bg': MSO_LANGUAGE_ID.BULGARIAN,
    'bn': MSO_LANGUAGE_ID.BENGALI,
    'bs': MSO_LANGUAGE_ID.BOSNIAN,
    'cs': MSO_LANGUAGE_ID.CZECH,
    'da': MSO_LANGUAGE_ID.DANISH,
    'de': MSO_LANGUAGE_ID.GERMAN,
    'el': MSO_LANGUAGE_ID.GREEK,
    'en': MSO_LANGUAGE_ID.ENGLISH_US,
    'es': MSO_LANGUAGE_ID.SPANISH,
    'et': MSO_LANGUAGE_ID.ESTONIAN,
    'fi': MSO_LANGUAGE_ID.FINNISH,
    'fr': MSO_LANGUAGE_ID.FRENCH,
    'fr-CA': MSO_LANGUAGE_ID.FRENCH_CANADIAN,
    'ha': MSO_LANGUAGE_ID.HAUSA,
    'he': MSO_LANGUAGE_ID.HEBREW,
    'hi': MSO_LANGUAGE_ID.HINDI,
    'hr': MSO_LANGUAGE_ID.CROATIAN,
    'hu': MSO_LANGUAGE_ID.HUNGARIAN,
    'id': MSO_LANGUAGE_ID.INDONESIAN,
    'it': MSO_LANGUAGE_ID.ITALIAN,
    'ja': MSO_LANGUAGE_ID.JAPANESE,
    'ka': MSO_LANGUAGE_ID.GEORGIAN,
    'ko': MSO_LANGUAGE_ID.KOREAN,
    'lv': MSO_LANGUAGE_ID.LATVIAN,
    'ms': MSO_LANGUAGE_ID.MALAYSIAN,
    'nl': MSO_LANGUAGE_ID.DUTCH,
    'no': MSO_LANGUAGE_ID.NORWEGIAN_BOKMOL,
    'pl': MSO_LANGUAGE_ID.POLISH,
    'ps': MSO_LANGUAGE_ID.PASHTO,
    'pt': MSO_LANGUAGE_ID.BRAZILIAN_PORTUGUESE,
    'ro': MSO_LANGUAGE_ID.ROMANIAN,
    'ru': MSO_LANGUAGE_ID.RUSSIAN,
    'sk': MSO_LANGUAGE_ID.SLOVAK,
    'sl': MSO_LANGUAGE_ID.SLOVENIAN,
    'so': MSO_LANGUAGE_ID.SOMALI,
    'sq': MSO_LANGUAGE_ID.ALBANIAN,
    'sr': MSO_LANGUAGE_ID.SERBIAN_LATIN,
    'sv': MSO_LANGUAGE_ID.SWEDISH,
    'sw': MSO_LANGUAGE_ID.SWAHILI,
    'ta': MSO_LANGUAGE_ID.TAMIL,
    'th': MSO_LANGUAGE_ID.THAI,
    'tr': MSO_LANGUAGE_ID.TURKISH,
    'uk': MSO_LANGUAGE_ID.UKRAINIAN,
    'ur': MSO_LANGUAGE_ID.URDU,
    'vi': MSO_LANGUAGE_ID.VIETNAMESE,
    'zh': MSO_LANGUAGE_ID.CHINESE_SINGAPORE,
    'zh-TW': MSO_LANGUAGE_ID.CHINESE_HONG_KONG_SAR,
}

TERMINOLOGY_NAME = 'pptx-translator-terminology'
skip = ''  # Text to skip

translate = boto3.client(service_name='translate', region_name='us-west-2',
                         aws_access_key_id='',
                         aws_secret_access_key='')


def resize(shape):
    # for size in [15, 13, 8, 2]:
    #   try:
    #      shape.text_frame.fit_text(max_size=size)
    #     break
    # except TypeError:
    #   pass
    # print('Could not fit text!')
    try:
        shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        shape.text_frame.word_wrap = True
    except AttributeError:
        traceback.print_exc()


def translate_request(text, source_language_code, dest_language_code, terminology_names):
    if len(text) > 0 and text.casefold() != skip.casefold() and not text.startswith("http"):
        try:
            response = translate.translate_text(
                Text=text,
                SourceLanguageCode=source_language_code,
                TargetLanguageCode=dest_language_code,
                TerminologyNames=terminology_names)
            return response.get('TranslatedText')
        except ClientError as client_error:
            if client_error.response['Error']['Code'] == 'ValidationException':
                # Text not valid. Maybe the size of the text exceeds the size limit of the service.
                # Amazon Translate limits: https://docs.aws.amazon.com/translate/latest/dg/what-is-limits.html
                # We just ignore and don't translate the text.
                print('Invalid text. Ignoring...')


def translate_main(slide, presentation, source_language_code, dest_language_code, terminology_names):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            paraText = ''
            for index, paragraph_run in enumerate(paragraph.runs):  # concatenate all the runs
                paraText += paragraph_run.text
            run = paragraph.add_run()
            # set font for translated text
            if paragraph.font.size is not None:
                run.font.size = paragraph.font.size - Pt(3)
                paragraph.font.size = run.font.size
            elif paragraph.runs[0].font.size is not None:
                run.font.size = paragraph.runs[0].font.size - Pt(3)
                for each in paragraph.runs: each.font.size = run.font.size
            # #Translate text with large enough font if size is not none
            if run.font.size is not None and run.font.size > Pt(13) or run.font.size is None:
                returned = translate_request(paraText, source_language_code, dest_language_code, terminology_names)
                if isinstance(returned, str):
                    run.text = " " + returned
            run.font.language_id = LANGUAGE_CODE_TO_LANGUAGE_ID[dest_language_code]
            # paragraph.runs[index].font.language_id = LANGUAGE_CODE_TO_LANGUAGE_ID[dest_language_code]
        resize(shape)


def translate_presentation(presentation, source_language_code, dest_language_code, terminology_names):
    slide_number = 1
    for slide in presentation.slides:
        print('Slide {slide_number} of {number_of_slides}'.format(
            slide_number=slide_number,
            number_of_slides=len(presentation.slides)))
        slide_number += 1
        translate_main(slide, presentation, source_language_code, dest_language_code, terminology_names)


def import_terminology(terminology_file_path):
    print('Importing terminology data from {file_path}...'.format(file_path=terminology_file_path))
    with open(terminology_file_path, 'rb') as f:
        translate.import_terminology(Name=TERMINOLOGY_NAME,
                                     MergeStrategy='OVERWRITE',
                                     TerminologyData={'File': bytearray(f.read()), 'Format': 'CSV'})


def read_skip_text():
    dir_path = os.path.dirname(os.path.realpath(__file__))
    txt_path = join(dir_path, 'workspace', 'skip.txt')
    try:
        text_input = open(txt_path)
        text_to_skip = text_input.read()
        text_input.close()
        global skip
        skip = text_to_skip
    except OSError:
        print('could not read file:' + txt_path)


def iterate_files(terminology_names, source_lang, dest_lang):
    read_skip_text()
    dir_path = os.path.dirname(os.path.realpath(__file__))
    input_dir = join(dir_path, 'workspace')
    file_paths = [join(input_dir, each) for each in listdir(input_dir) if isfile(join(input_dir, each))
                  and each.endswith('.pptx')]
    out_dir = join(input_dir, 'output')
    Path(out_dir).mkdir(parents=True, exist_ok=True)

    for file_path in file_paths:
        print('Translating {file_path} from {source_language_code} to {dest_language_code}...'.format(
            file_path=file_path,
            source_language_code=source_lang,
            dest_language_code=dest_lang))
        presentation = Presentation(file_path)
        translate_presentation(presentation, source_lang, dest_lang, terminology_names)
        output_file_path = join(input_dir + '\\' + 'output\\' + source_lang + '-' + dest_lang + '-' + date.today().strftime('%m-%d-') + os.path.basename(file_path))
        print('Saving {output_file_path}...'.format(output_file_path=output_file_path))
        presentation.save(output_file_path)

def main():
    argument_parser = argparse.ArgumentParser(
        'Translates pptx files from source language to target language using Amazon Translate service')
    argument_parser.add_argument(
        '-t', type=str,
        help='The path of the terminology CSV file')
    argument_parser.add_argument(
        '-s', '--source_lang', type=str,
        help='The source language')
    argument_parser.add_argument(
        '-d', '--dest_lang', type=str,
        help='The destination language')
    args = argument_parser.parse_args()

    terminology_names = []
    if args.t:
        import_terminology(args.t)
        terminology_names = [TERMINOLOGY_NAME]
    source_lang = args.source_lang
    dest_lang = args.dest_lang
    iterate_files(terminology_names, source_lang, dest_lang)

if __name__ == '__main__':
    main()
    input("Press enter to exit;")
